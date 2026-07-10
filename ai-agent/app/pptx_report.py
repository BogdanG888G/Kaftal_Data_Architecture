"""Генерация PPTX-презентации из отчёта."""
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from charts import fig_to_png_bytes


# Цвета темы (совпадают с UI)
BG_DARK_RGB = (13, 17, 23)
CARD_BG_RGB = (22, 27, 34)
TEXT_MAIN_RGB = (240, 246, 252)
TEXT_SUB_RGB = (139, 148, 158)
ACCENT_RED = (255, 107, 107)
ACCENT_TEAL = (0, 210, 255)
ACCENT_YELLOW = (255, 217, 61)
ACCENT_PURPLE = (192, 132, 252)
ACCENT_ORANGE = (255, 159, 67)
ACCENT_GREEN = (107, 203, 119)
ACCENT_BLUE = (77, 157, 224)
ACCENT_PINK = (255, 110, 180)

KPI_COLORS = [
    ACCENT_RED, ACCENT_TEAL, ACCENT_YELLOW, ACCENT_PURPLE,
    ACCENT_ORANGE, ACCENT_GREEN, ACCENT_BLUE, ACCENT_PINK,
]


# Размеры слайда: широкий формат 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# Русские названия для метрик и фильтров
LABEL_TRANSLATIONS = {
    # Метрики
    "revenue": "Выручка",
    "qty": "Штук",
    "tt_count": "Торговых точек",
    "stores": "Магазинов",
    "brands": "Брендов",
    "brands_count": "Брендов",
    "flavors": "Вкусов",
    "flavors_count": "Вкусов",
    "products": "Товаров",
    "products_count": "Товаров",
    "avg_price": "Средняя цена",
    "avg_cost": "Средняя себестоимость",
    "avg_cost_price": "Средняя себестоимость",
    "margin": "Маржа",
    "margin_rub": "Маржа",
    "margin_pct": "Маржа %",

    # Top-N
    "top_brand": "Топ-1 бренд",
    "top_flavor": "Топ-1 вкус",
    "top_chain": "Топ-1 сеть",
    "top_region": "Топ-1 регион",
    "top_city": "Топ-1 город",

    # Фильтры
    "retail_chain": "Сеть",
    "store_format": "Формат магазина",
    "chip_type": "Тип чипсов",
    "year": "Год",
    "month": "Месяц",
    "weight_grams": "Граммовка",
    "weight_grams_list": "Граммовки",
    "priority_flavors": "Приоритетные вкусы",
    "brand": "Бренд",
    "flavor": "Вкус",
}


def _clone_fig_without_title(fig):
    """Возвращает копию фигуры без title (для PPTX)."""
    import plotly.graph_objects as go

    fig_dict = fig.to_dict()

    if "layout" not in fig_dict:
        fig_dict["layout"] = {}

    fig_dict["layout"]["title"] = {"text": ""}

    if "margin" not in fig_dict["layout"]:
        fig_dict["layout"]["margin"] = {}
    fig_dict["layout"]["margin"]["t"] = 20

    return go.Figure(fig_dict)


def translate(name: str) -> str:
    return LABEL_TRANSLATIONS.get(name.lower(), name)


def _should_skip_kpi(col_name: str, value) -> bool:
    """
    Определяет, нужно ли пропустить бесполезную KPI-карточку.
    Например: tt_count=1, stores=1 и т.п.
    """
    import pandas as pd

    try:
        if value is None or pd.isna(value):
            return True
    except (TypeError, ValueError):
        pass

    col_lower = col_name.lower()

    # Для количественных единичных значений — скрываем если == 1
    skip_if_one = ("tt_count", "stores", "stores_count",
                   "brands", "brands_count",
                   "flavors", "flavors_count",
                   "products", "products_count")

    if col_lower in skip_if_one:
        try:
            if float(value) <= 1:
                return True
        except (ValueError, TypeError):
            pass

    # Пустые строковые значения
    if isinstance(value, str) and value.strip() in ("", "Не указано", "None"):
        return True

    return False


def fmt_value(val, col_name=""):
    """Форматирование значения для карточки."""
    import pandas as pd

    if isinstance(val, (list, tuple)):
        if len(val) == 0:
            return "—"
        items = [str(v) for v in val[:3]]
        result = ", ".join(items)
        if len(val) > 3:
            result += f" +{len(val) - 3}"
        return result

    try:
        if val is None or pd.isna(val):
            return "—"
    except (TypeError, ValueError):
        pass

    col_lower = col_name.lower()

    if any(k in col_lower for k in ["revenue", "rub", "amount", "cost", "margin_rub", "price"]) and "pct" not in col_lower:
        try:
            n = float(val)
            if n >= 1_000_000_000:
                return f"{n/1_000_000_000:.2f} млрд ₽"
            if n >= 1_000_000:
                return f"{n/1_000_000:.2f} млн ₽"
            if n >= 1_000:
                return f"{n/1_000:.1f} тыс ₽"
            return f"{n:.2f} ₽"
        except (ValueError, TypeError):
            return str(val)

    if "pct" in col_lower:
        try:
            return f"{float(val):.1f}%"
        except (ValueError, TypeError):
            return str(val)

    if any(k in col_lower for k in ["qty", "quantity", "count", "stores", "brands", "flavors", "products"]):
        try:
            n = float(val)
            if n >= 1_000_000:
                return f"{n/1_000_000:.1f} млн"
            if n >= 1_000:
                return f"{n:,.0f}".replace(",", " ")
            return f"{int(n)}"
        except (ValueError, TypeError):
            return str(val)

    if isinstance(val, (int, float)):
        try:
            if val >= 1000:
                return f"{val:,.2f}".replace(",", " ")
            return f"{val:.2f}" if isinstance(val, float) else str(val)
        except Exception:
            return str(val)

    return str(val)


def _set_bg_color(slide, rgb):
    """Устанавливает цвет фона слайда."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _add_text(slide, text, left, top, width, height,
              size=18, bold=False, color=TEXT_MAIN_RGB,
              align=PP_ALIGN.LEFT):
    """Добавляет текстовый блок."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Arial"
    return txbox


def _add_kpi_card(slide, left, top, width, height, label, value, color):
    """Рисует KPI-карточку."""
    # Основной прямоугольник
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*CARD_BG_RGB)
    card.line.color.rgb = RGBColor(*color)
    card.line.width = Pt(2)

    # Верхняя цветная полоса
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Emu(50000), top + Emu(30000),
        width - Emu(100000), Emu(80000),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(*color)
    top_bar.line.fill.background()

    # Метка (label)
    label_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.3),
        width - Inches(0.4), Inches(0.4),
    )
    lf = label_box.text_frame
    lp = lf.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lrun = lp.add_run()
    lrun.text = label.upper()
    lrun.font.size = Pt(11)
    lrun.font.bold = True
    lrun.font.color.rgb = RGBColor(*TEXT_SUB_RGB)
    lrun.font.name = "Arial"

    # Значение
    val_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.7),
        width - Inches(0.4), Inches(0.8),
    )
    vf = val_box.text_frame
    vp = vf.paragraphs[0]
    vp.alignment = PP_ALIGN.LEFT
    vrun = vp.add_run()
    vrun.text = value
    vrun.font.size = Pt(22)
    vrun.font.bold = True
    vrun.font.color.rgb = RGBColor(*color)
    vrun.font.name = "Arial"


def build_pptx_bytes(report: dict) -> bytes:
    """
    Генерирует PPTX-презентацию из отчёта.

    Args:
        report: dict от build_report()

    Returns:
        bytes для скачивания
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # blank

    # ==================================================
    # СЛАЙД 1: ТИТУЛЬНЫЙ
    # Только заголовок отчёта + применённые фильтры
    # ==================================================
    slide = prs.slides.add_slide(blank_layout)
    _set_bg_color(slide, BG_DARK_RGB)

    # Заголовок — сам запрос пользователя, крупно по центру
    _add_text(
        slide,
        report["request"],
        Inches(1), Inches(2.5),
        Inches(11), Inches(2.5),
        size=28, bold=True, color=TEXT_MAIN_RGB, align=PP_ALIGN.CENTER,
    )

    # Применённые фильтры под заголовком
    filters = report.get("filters") or {}
    active = {k: v for k, v in filters.items() if v is not None and v != []}
    if active:
        filter_lines = []
        for k, v in list(active.items())[:8]:
            label = translate(k)
            if isinstance(v, list):
                val_str = ", ".join(str(x) for x in v)
            else:
                val_str = str(v)
            filter_lines.append(f"• {label}: {val_str}")
        _add_text(
            slide,
            "\n".join(filter_lines),
            Inches(2), Inches(5.5),
            Inches(9), Inches(2),
            size=14, color=TEXT_SUB_RGB, align=PP_ALIGN.CENTER,
        )

    # ==================================================
    # СЛАЙДЫ ПО СЕКЦИЯМ
    # ==================================================
    for section in report["sections"]:
        if section.get("error"):
            continue

        title = section["title"]
        chart_type = section.get("chart_type")
        df = section.get("data")
        fig = section.get("fig")

        # KPI-слайд особый
        if chart_type == "kpi" and df is not None and len(df) == 1:
            slide = prs.slides.add_slide(blank_layout)
            _set_bg_color(slide, BG_DARK_RGB)

            _add_text(
                slide, title.upper(),
                Inches(0.5), Inches(0.3),
                Inches(12), Inches(0.7),
                size=24, bold=True, color=TEXT_MAIN_RGB,
            )

            # Фильтруем колонки: пропускаем бесполезные (tt_count=1, пустые строки)
            row = df.iloc[0]
            cols = [c for c in df.columns if not _should_skip_kpi(c, row[c])][:8]

            if not cols:
                _add_text(
                    slide, "Нет значимых KPI для отображения",
                    Inches(1), Inches(3),
                    Inches(11), Inches(1),
                    size=16, color=TEXT_SUB_RGB, align=PP_ALIGN.CENTER,
                )
                continue

            per_row = 4
            card_w = Inches(2.9)
            card_h = Inches(1.6)
            gap_x = Inches(0.2)
            gap_y = Inches(0.3)
            start_x = Inches(0.5)
            start_y = Inches(1.5)

            for i, col_name in enumerate(cols):
                r = i // per_row
                c = i % per_row
                x = start_x + c * (card_w + gap_x)
                y = start_y + r * (card_h + gap_y)
                val = row[col_name]
                _add_kpi_card(
                    slide, x, y, card_w, card_h,
                    label=translate(col_name),
                    value=fmt_value(val, col_name),
                    color=KPI_COLORS[i % len(KPI_COLORS)],
                )

            continue

        # Обычный слайд с графиком
        if fig is None:
            continue

        slide = prs.slides.add_slide(blank_layout)
        _set_bg_color(slide, BG_DARK_RGB)

        # Заголовок слайда
        _add_text(
            slide, title.upper(),
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7),
            size=22, bold=True, color=TEXT_MAIN_RGB,
        )

        # График (без заголовка внутри — чтобы не дублировать)
        try:
            fig_no_title = _clone_fig_without_title(fig)
            png_bytes = fig_to_png_bytes(fig_no_title, width=1600, height=900)
            if png_bytes:
                img_stream = io.BytesIO(png_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0.5), Inches(1.2),
                    Inches(12.3), Inches(6.0),
                )
        except Exception as e:
            print(f"[PPTX] Failed to add chart for {title}: {e}")
            _add_text(
                slide, f"[График недоступен: {e}]",
                Inches(1), Inches(3),
                Inches(11), Inches(1),
                size=14, color=ACCENT_RED,
            )

    # ==================================================
    # СОХРАНЯЕМ
    # (Финальный слайд "СПАСИБО" убран по требованию)
    # ==================================================
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()